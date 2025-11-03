#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PPT ingester for scout.impact_repository

- Scans a folder for .pptx/.ppt decks
- Extracts slide text (per slide)
- Classifies impact pointers (LLM) -> (Industry, Business Group, Use Case) + confidences
- Generates text embeddings
- Deduplicates via SHA-256 fingerprint
- Upserts into Postgres table `scout.impact_repository` (ON CONFLICT DO NOTHING on fingerprint)

Usage:
  python pipelines/run_ppt_ingest.py \
      --folder "C:\\...\\Stakeholder_360\\data\\case_studies" \
      --min-conf 0.80 \
      --top-n 10 \
      --concurrency 6
"""

import os
import re
import json
import time
import math
import hashlib
import argparse
import asyncio
from typing import Any, Dict, List, Optional, Tuple

import nest_asyncio
nest_asyncio.apply()

import pandas as pd
from tqdm import tqdm
from pptx import Presentation

from dotenv import load_dotenv
load_dotenv()

# --- Azure OpenAI (Chat) ---
from openai import AsyncAzureOpenAI

# --- DB / SQLAlchemy ---
from sqlalchemy import MetaData, Table, text
from sqlalchemy.dialects.postgresql import insert as pg_insert
from sqlalchemy.types import UserDefinedType

# --- Local utils ---
from utils.db import get_engine
from pipelines.embedding_utils import generate_embeddings, EMBEDDING_DIM

# =============================================================================
# Register pgvector type with SQLAlchemy reflection
# =============================================================================
class Vector(UserDefinedType):
    """
    Minimal SQLAlchemy type for pgvector columns.
    Accepts arbitrary args/kwargs so reflection doesn't crash (vector(1536)).
    """
    def __init__(self, *args, **kwargs):
        pass
    def get_col_spec(self) -> str:
        # Return generic spec; actual dim maintained in DB. Not used for reflection.
        return "vector"

try:
    from sqlalchemy.dialects.postgresql.base import ischema_names
    ischema_names["vector"] = Vector  # type: ignore[index]
except Exception:
    # Older SA fallback
    from sqlalchemy.dialects import postgresql as _pg
    _pg.base.ischema_names["vector"] = Vector  # type: ignore[attr-defined]

# =============================================================================
# Constants / Env
# =============================================================================
TABLE_SCHEMA = "scout"
TABLE_NAME = "impact_repository"
FQN = f"{TABLE_SCHEMA}.{TABLE_NAME}"

DEFAULT_MIN_CONF = float(os.getenv("MIN_CONF", "0.80"))
DEFAULT_TOP_N = int(os.getenv("TOP_N_PER_DECK", "10"))
DEFAULT_CONCURRENCY = int(os.getenv("CLASSIFY_CONCURRENCY", "6"))

# Azure OpenAI env
AZURE_ENDPOINT    = os.getenv("AZURE_ENDPOINT", "")
AZURE_DEPLOYMENT  = os.getenv("AZURE_DEPLOYMENT", "")
AZURE_API_KEY     = os.getenv("AZURE_API_KEY", "")
AZURE_API_VERSION = os.getenv("AZURE_API_VERSION", "2024-02-15-preview")

# Classification prompts
CLASSIFY_PROMPT = """
You are a domain expert in extracting business impact pointers from slide text and classifying them.

TASK
From the provided slide text, first **extract impact pointers verbatim** (no paraphrasing). An impact pointer is a measurable outcome line
(e.g., contains a number, %, $, hours, accuracy, savings, lift, reduction/increase, time, cost, revenue, throughput, etc.). Keep punctuation and wording EXACTLY as in the slide.

Then, for EACH extracted impact pointer, classify into:
1. Industry — macro domain (CPG, Retail, Manufacturing, Automotive, BFSI, Pharma, Healthcare, Utilities, Telecom, Technology, Public Sector, Transportation/Logistics, etc.)
2. Business Group — function (Marketing, Sales, Finance, Procurement, Supply Chain, Operations, Manufacturing, Product/Engineering, IT, Executive/CEO Office, Commercial Analytics, Revenue Management, etc.)
3. Use Case — short, precise title (e.g., "GenAI Contract Data Extraction for Billing Accuracy", "Executive Dashboard for Boardroom Insights", "Scenario Management for Production Scheduling", "Referral Journey Analytics").

Confidence scores must be between 0 and 1 in **0.05 steps only** (0.5, 0.55, 0.6, 0.65, 0.7, 0.75, 0.8, 0.85, 0.9, 0.95, 1.0). If multiple options exist, include up to 3 ranked options per field.

OUTPUT
Return ONLY valid JSON:
1. Industry — macro business domain (CPG, Retail, Manufacturing, BFSI, Pharma, Utilities, etc.)
2. Business Group — function (Marketing, Sales, Supply Chain, Finance, Manufacturing, etc.)
3. Use Case — concise description of the analytical or business initiative.

### Confidence Logic
For each classification, include a numeric confidence score between **0 and 1**, but restrict it to clean 0.05 intervals only
(e.g., 0.5, 0.55, 0.6, 0.65, 0.7, 0.75, 0.8, 0.85, 0.9, 0.95, 1.0).  
If multiple options exist, include up to 3 ranked options per field, sorted by descending confidence.

Return **only valid JSON** in this structure:
{{
  "impacts": [
    {{
      "text": "<VERBATIM impact pointer from slides>",
      "Industry": [{{"value": "CPG", "confidence": 0.95}}],
      "Business Group": [{{"value": "Finance & Revenue Management", "confidence": 0.9}}],
      "Use Case": [{{"value": "GenAI Contract Data Extraction for Billing Accuracy", "confidence": 0.95}}]
    }}
  ]
}}

Rules:
- If no impact pointers are present, return {{ "impacts": [] }}.
- Preserve pointers EXACTLY as in slides (no summarization).

Slide Text:
{slide_text}
""".strip()

CLEANUP_PROMPT = """
You are an expert data cleaner for impact statements.

Input is a JSON object with a key "impacts".
Each element has "text", "Industry", "Business Group", and "Use Case".

✅ Keep only entries describing meaningful business outcomes 
(e.g., savings, efficiency, automation, accuracy, revenue, cost, hours, improvement).

🚫 Remove irrelevant numeric-only, raw, or incomplete entries
(e.g., "Unit Shortfall: 1,317", "Damages Owed: $395,100").

Return only valid JSON in the same structure, no commentary.

Input JSON:
{raw_json}
""".strip()

# =============================================================================
# Helpers
# =============================================================================

def compute_fingerprint(row: Dict[str, Any]) -> str:
    """
    Deterministic unique hash for dedupe.
    """
    sig = "|".join([
        (row.get("impact") or "").strip().lower(),
        (row.get("industry") or "").strip().lower(),
        (row.get("business_group") or "").strip().lower(),
        (row.get("use_case") or "").strip().lower(),
        (row.get("source_file") or "").strip().lower(),
        (row.get("source_type") or "").strip().lower(),
    ])
    return hashlib.sha256(sig.encode("utf-8")).hexdigest()

def safe_json_load(s: str) -> Dict[str, Any]:
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        return {"impacts": []}

def expand_impacts_to_rows(
    data: Dict[str, Any],
    source_file: str,
    slide_no: Optional[int]
) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    for imp in data.get("impacts", []):
        impact_text = (imp.get("text") or "").strip()
        industries = imp.get("Industry") or [{"value": "", "confidence": 0}]
        bgroups   = imp.get("Business Group") or [{"value": "", "confidence": 0}]
        usecases  = imp.get("Use Case") or [{"value": "", "confidence": 0}]
        for ind in industries:
            for bg in bgroups:
                for uc in usecases:
                    rows.append({
                        "source_file": source_file,
                        "impact": impact_text,
                        "industry": (ind.get("value") or "").strip(),
                        "industry_conf": float(ind.get("confidence") or 0),
                        "business_group": (bg.get("value") or "").strip(),
                        "business_group_conf": float(bg.get("confidence") or 0),
                        "use_case": (uc.get("value") or "").strip(),
                        "use_case_conf": float(uc.get("confidence") or 0),
                        "source_type": "ppt",
                        "slide_no": slide_no,
                        "embedding": None,    # fill later
                        "fingerprint": "",    # fill later
                    })
    return rows

def filter_confident(df: pd.DataFrame, min_conf: float) -> pd.DataFrame:
    if df.empty:
        return df
    out = df[
        (pd.to_numeric(df["industry_conf"], errors="coerce")       >= min_conf) &
        (pd.to_numeric(df["business_group_conf"], errors="coerce") >= min_conf) &
        (pd.to_numeric(df["use_case_conf"], errors="coerce")       >= min_conf)
    ].copy()
    return out

# =============================================================================
# Azure OpenAI client + classify with retries
# =============================================================================

def build_azure_client() -> AsyncAzureOpenAI:
    if not (AZURE_ENDPOINT and AZURE_DEPLOYMENT and AZURE_API_KEY):
        raise RuntimeError("Azure OpenAI configuration missing. Set AZURE_ENDPOINT, AZURE_DEPLOYMENT, AZURE_API_KEY.")
    return AsyncAzureOpenAI(
        azure_endpoint=AZURE_ENDPOINT,
        api_key=AZURE_API_KEY,
        api_version=AZURE_API_VERSION,
        timeout=90.0
    )

async def _chat_json(client: AsyncAzureOpenAI, user_content: str, max_retries: int = 3, retry_backoff: float = 2.0) -> str:
    last_err = None
    for attempt in range(1, max_retries + 1):
        try:
            resp = await client.chat.completions.create(
                model=AZURE_DEPLOYMENT,
                temperature=0.0,
                messages=[
                    {"role": "system", "content": "Return only valid JSON. No explanations."},
                    {"role": "user", "content": user_content},
                ],
            )
            reply = (resp.choices[0].message.content or "").strip()
            # strip accidental code fences
            reply = reply.replace("```json", "").replace("```", "").strip()
            return reply
        except Exception as e:
            last_err = e
            if attempt < max_retries:
                await asyncio.sleep(retry_backoff * attempt)
            else:
                return '{"impacts": []}'
    raise last_err or RuntimeError("Unknown Azure OpenAI error")

async def classify_slide_text(client: AsyncAzureOpenAI, slide_text: str) -> Dict[str, Any]:
    raw = await _chat_json(client, CLASSIFY_PROMPT.format(slide_text=slide_text))
    return safe_json_load(raw)

async def cleanup_impacts(client: AsyncAzureOpenAI, raw_json_str: str) -> Dict[str, Any]:
    cleaned = await _chat_json(client, CLEANUP_PROMPT.format(raw_json=raw_json_str))
    return safe_json_load(cleaned)

# =============================================================================
# DB ensure + upsert
# =============================================================================

DDL_ALTER_ADD_VECTOR = f"""
ALTER TABLE {FQN}
    ADD COLUMN IF NOT EXISTS embedding vector({EMBEDDING_DIM});
"""

def ensure_table_ready(engine) -> None:
    """
    Ensure the target table has `embedding vector(dim)`.
    """
    with engine.begin() as conn:
        conn.execute(text(DDL_ALTER_ADD_VECTOR))

def upsert_dataframe(engine, df: pd.DataFrame) -> int:
    """
    Upsert rows into scout.impact_repository using ON CONFLICT DO NOTHING on fingerprint.
    Drops any extra columns not in table.
    """
    if df.empty:
        return 0

    metadata = MetaData(schema=TABLE_SCHEMA)
    table = Table(TABLE_NAME, metadata, autoload_with=engine)

    valid_cols = set(table.columns.keys())
    df2 = df[[c for c in df.columns if c in valid_cols]].copy()

    records = df2.to_dict(orient="records")

    with engine.begin() as conn:
        stmt = pg_insert(table).values(records)
        stmt = stmt.on_conflict_do_nothing(index_elements=["fingerprint"])
        res = conn.execute(stmt)
        return res.rowcount or 0

# =============================================================================
# PPT parsing
# =============================================================================

def extract_slide_texts(ppt_path: str) -> List[Tuple[int, str]]:
    """
    Return list of (slide_no, text) for a PPTX/PPT.
    Basic text extraction from shapes; tables/notes not included here by design.
    """
    prs = Presentation(ppt_path)
    out: List[Tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        blocks = [
            shp.text.strip()
            for shp in slide.shapes
            if hasattr(shp, "text") and isinstance(shp.text, str) and shp.text.strip()
        ]
        if blocks:
            out.append((i, "\n".join(blocks)))
    return out

# =============================================================================
# Main
# =============================================================================

async def main(
    folder: str,
    min_conf: float,
    top_n: int,
    concurrency: int
) -> None:
    if not os.path.isdir(folder):
        raise FileNotFoundError(f"Folder not found: {folder}")

    # Collect PPT files
    ppt_files = [
        os.path.join(folder, f)
        for f in os.listdir(folder)
        if f.lower().endswith((".pptx", ".ppt"))
    ]
    if not ppt_files:
        print(f"ℹ️ No PPT/PPTX found under: {folder}")
        return

    client = build_azure_client()

    all_rows: List[Dict[str, Any]] = []

    for ppt_path in tqdm(ppt_files, desc="Scanning decks"):
        source_file = os.path.basename(ppt_path)
        slides = extract_slide_texts(ppt_path)
        if not slides:
            print(f"🟡 No text extracted from {source_file}")
            continue

        # Classify slides with bounded concurrency
        sem = asyncio.Semaphore(concurrency)

        async def _proc_slide(slide_no: int, text: str) -> List[Dict[str, Any]]:
            async with sem:
                raw = await classify_slide_text(client, text)
                # optional cleanup pass to drop numeric-only lines
                cleaned = await cleanup_impacts(client, json.dumps(raw))
                rows = expand_impacts_to_rows(cleaned, source_file=source_file, slide_no=slide_no)
                return rows

        tasks = [_proc_slide(sn, st) for (sn, st) in slides]
        per_slide_rows = await asyncio.gather(*tasks)

        # Flatten
        for chunk in per_slide_rows:
            all_rows.extend(chunk)

    # Build DataFrame
    df = pd.DataFrame(all_rows, columns=[
        "source_file","impact","industry","industry_conf",
        "business_group","business_group_conf","use_case","use_case_conf",
        "source_type","slide_no","embedding","fingerprint"
    ])

    if df.empty:
        print("\nℹ️ Nothing extracted. Exiting.")
        return

    # Confidence filter
    df = filter_confident(df, min_conf=min_conf)
    if df.empty:
        print(f"\nℹ️ After confidence filter (min_conf={min_conf}), nothing to ingest.")
        return

    # Average conf + Top-N per deck
    df["avg_conf"] = (
        pd.to_numeric(df["industry_conf"], errors="coerce").fillna(0) +
        pd.to_numeric(df["business_group_conf"], errors="coerce").fillna(0) +
        pd.to_numeric(df["use_case_conf"], errors="coerce").fillna(0)
    ) / 3.0

    if top_n > 0:
        df = df.sort_values(["source_file", "avg_conf"], ascending=[True, False])
        df = df.groupby("source_file", as_index=False).head(top_n)

    if df.empty:
        print(f"\nℹ️ After top-n={top_n}, nothing to ingest.")
        return

    # Fingerprints (before dedupe)
    df["fingerprint"] = df.apply(lambda r: compute_fingerprint(r.to_dict()), axis=1)
    df = df.drop_duplicates(subset=["fingerprint"]).copy()

    # Embeddings
    embed_texts = (df["impact"].fillna("") + " | " +
                   df["use_case"].fillna("") + " | " +
                   df["industry"].fillna("") + " | " +
                   df["business_group"].fillna("")).tolist()

    print(f"🧠 Generating embeddings for {len(embed_texts)} rows (dim={EMBEDDING_DIM})...")

    vectors = await generate_embeddings(df["impact"].tolist())

    
    if len(vectors) != len(df):
        raise RuntimeError("Embedding count mismatch.")
    df["embedding"] = vectors

    # Upsert to DB
    engine = get_engine()
    ensure_table_ready(engine)
    inserted = upsert_dataframe(engine, df)

    # Summary
    print("\n================== Summary ==================")
    print(f"Decks scanned   : {len(ppt_files)}")
    print(f"Rows extracted  : {len(all_rows)} (pre-filter)")
    print(f"Rows kept       : {len(df)} (after confidence, top-n, dedupe)")
    print(f"Inserted        : {inserted} (ON CONFLICT DO NOTHING on fingerprint)")
    print("Target table    :", FQN)
    print("=============================================\n")

# =============================================================================
# CLI
# =============================================================================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ingest PPT impacts into Postgres (scout.impact_repository)")
    parser.add_argument("--folder", required=True, help="Folder containing PPT/PPTX files")
    parser.add_argument("--min-conf", type=float, default=DEFAULT_MIN_CONF, help="Min confidence threshold (all 3 dims)")
    parser.add_argument("--top-n", type=int, default=DEFAULT_TOP_N, help="Top-N per deck (post-filter). 0=all")
    parser.add_argument("--concurrency", type=int, default=DEFAULT_CONCURRENCY, help="Max concurrent slide classifications")
    args = parser.parse_args()

    asyncio.run(main(
        folder=args.folder,
        min_conf=args.min_conf,
        top_n=args.top_n,
        concurrency=args.concurrency,
    ))
