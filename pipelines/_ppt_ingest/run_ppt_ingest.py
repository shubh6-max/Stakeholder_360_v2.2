import os
import re
import json
import hashlib
import nest_asyncio
import pandas as pd
import asyncio
from pptx import Presentation
from openai import AsyncAzureOpenAI
from tqdm import tqdm
from typing import List, Dict, Any, Optional
from sqlalchemy.dialects.postgresql import insert as pg_insert
from sqlalchemy import text
from sqlalchemy import MetaData, Table
from sqlalchemy.dialects.postgresql import insert as pg_insert
from dotenv import load_dotenv
from utils.db import get_engine
load_dotenv()
# =========================================
# ✅ Register pgvector column type cleanly
# =========================================
from sqlalchemy.types import UserDefinedType
from sqlalchemy.dialects import registry

class Vector(UserDefinedType):
    """Minimal SQLAlchemy type for pgvector columns (reflection-safe)."""
    def __init__(self, *args, **kwargs):
        # accept arbitrary args/kwargs for reflection
        pass

    def get_col_spec(self):
        return "vector"

    def bind_processor(self, dialect):
        return None

    def result_processor(self, dialect, coltype):
        return None



# =======================================
# 🔧 Azure OpenAI Configuration
# =======================================
AZURE_ENDPOINT = os.getenv("AZURE_ENDPOINT", "https://openai-scout-001.openai.azure.com/")
AZURE_DEPLOYMENT = os.getenv("AZURE_DEPLOYMENT", "gpt-4-vision-preview")
AZURE_API_KEY = os.getenv("AZURE_API_KEY", "REPLACE_ME")
AZURE_API_VERSION = os.getenv("AZURE_API_VERSION", "2024-02-15-preview")

client = AsyncAzureOpenAI(
    azure_endpoint=AZURE_ENDPOINT,
    api_key=AZURE_API_KEY,
    api_version=AZURE_API_VERSION,
)

nest_asyncio.apply()

# =======================================
# ⚙️ Ingest Settings
# =======================================
TABLE_NAME = "scout.impact_repository"
SOURCE_TYPE = "ppt"                    # fixed for this pipeline
TOP_N_PER_DECK = int(os.getenv("TOP_N_PER_DECK", "10"))   # change here if you want
MIN_CONF = float(os.getenv("MIN_CONF", "0.80"))           # used for all 3 conf cols
FOLDER_PATH = os.getenv("PPT_FOLDER", r"/data/case_studies")   # can override from env

# =======================================
# 🗃️ DDL Ensure (idempotent)
# =======================================
DDL_CREATE = f"""
CREATE SCHEMA IF NOT EXISTS scout;

CREATE TABLE IF NOT EXISTS {TABLE_NAME} (
  id                   BIGSERIAL PRIMARY KEY,
  source_file          TEXT NOT NULL,
  impact               TEXT NOT NULL,
  industry             TEXT,
  industry_conf        DOUBLE PRECISION,
  business_group       TEXT,
  business_group_conf  DOUBLE PRECISION,
  use_case             TEXT,
  use_case_conf        DOUBLE PRECISION,
  source_type          TEXT NOT NULL DEFAULT 'ppt',
  slide_no             INTEGER,
  created_at           TIMESTAMPTZ DEFAULT NOW(),
  fingerprint          TEXT NOT NULL UNIQUE
);
"""

# =======================================
# 🔤 PPT Text Extraction (basic version you pasted)
# (If you want the robust extractor that covers tables & notes, I can drop it in.)
# =======================================
def extract_text_from_ppt(ppt_path: str) -> str:
    prs = Presentation(ppt_path)
    slide_texts = []
    for i, slide in enumerate(prs.slides, start=1):
        text_blocks = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if text_blocks:
            slide_texts.append(f"--- Slide {i} ---\n" + "\n".join(text_blocks))
    return "\n\n".join(slide_texts)

# =======================================
# 🤖 LLM Calls
# =======================================
async def analyze_ppt_text(ppt_text: str) -> str:
    prompt = """
You are a domain expert in extracting business impact pointers from slide text and classifying them.

TASK
From the provided slide text, first **extract impact pointers verbatim** (no paraphrasing). An impact pointer is a measurable outcome line
(e.g., contains a number, %, $, hours, accuracy, savings, lift, reduction/increase, time, cost, revenue, throughput, etc.). Keep punctuation and wording EXACTLY as in the slide.

Then, for EACH extracted impact pointer, classify into:
1. Industry — macro domain (CPG, Retail, Manufacturing, Automotive, BFSI, Pharma, Healthcare, Utilities, Telecom, Technology, Public Sector, Transportation/Logistics, etc.)
2. Business Group — function (Marketing, Sales, Finance, Procurement, Supply Chain, Operations, Manufacturing, Product/Engineering, IT, Executive/CEO Office, Commercial Analytics, Revenue Management, etc.)
3. Use Case — short, precise title (e.g., "GenAI Contract Data Extraction for Billing Accuracy", "Executive Dashboard for Boardroom Insights", "Scenario Management for Production Scheduling", "Referral Journey Analytics").

Confidence scores must be between 0 and 1 in **0.05 steps only** (0.5, 0.55, 0.6, 0.65, 0.7, 0.75, 0.8, 0.85, 0.9, 0.95, 1.0). If multiple options exist, include up to 3 ranked options per field.

OUTPUT
Return ONLY valid JSON:
1. Industry — macro business domain (CPG, Retail, Manufacturing, BFSI, Pharma, Utilities, etc.)
2. Business Group — function (Marketing, Sales, Supply Chain, Finance, Manufacturing, etc.)
3. Use Case — concise description of the analytical or business initiative.

### Confidence Logic
For each classification, include a numeric confidence score between **0 and 1**, but restrict it to clean 0.05 intervals only
(e.g., 0.5, 0.55, 0.6, 0.65, 0.7, 0.75, 0.8, 0.85, 0.9, 0.95, 1.0).  
If multiple options exist, include up to 3 ranked options per field, sorted by descending confidence.

Return **only valid JSON** in this structure:
{{
  "impacts": [
    {{
      "text": "<VERBATIM impact pointer from slides>",
      "Industry": [{{"value": "CPG", "confidence": 0.95}}],
      "Business Group": [{{"value": "Finance & Revenue Management", "confidence": 0.9}}],
      "Use Case": [{{"value": "GenAI Contract Data Extraction for Billing Accuracy", "confidence": 0.95}}]
    }}
  ]
}}

Rules:
- If no impact pointers are present, return {{ "impacts": [] }}.
- Preserve pointers EXACTLY as in slides (no summarization).

Slide Text:
{{ppt_text}}
""".strip()

    try:
        response = await client.chat.completions.create(
            model=AZURE_DEPLOYMENT,
            temperature=0,
            messages=[
                {"role": "system", "content": "Return only valid JSON, no explanations."},
                {"role": "user", "content": prompt.replace("{ppt_text}", ppt_text)},
            ],
        )
        reply = response.choices[0].message.content or ""
        return reply.replace("```json", "").replace("```", "").strip()
    except Exception as e:
        print(f"⚠️ LLM error: {e}")
        return '{"impacts": []}'

async def cleanup_irrelevant_impacts(raw_json_str: str) -> str:
    prompt_cleanup = f"""
You are an expert data cleaner for impact statements.

Input is a JSON object with a key "impacts".
Each element has "text", "Industry", "Business Group", and "Use Case".

✅ Keep only entries describing meaningful business outcomes 
(e.g., savings, efficiency, automation, accuracy, revenue, cost, hours, improvement).

🚫 Remove irrelevant numeric-only, raw, or incomplete entries
(e.g., "Unit Shortfall: 1,317", "Damages Owed: $395,100").

Return only valid JSON in the same format, no commentary.

Input JSON:
{raw_json_str}
"""
    try:
        response = await client.chat.completions.create(
            model=AZURE_DEPLOYMENT,
            temperature=0,
            messages=[
                {"role": "system", "content": "Return only valid JSON, no explanations."},
                {"role": "user", "content": prompt_cleanup},
            ],
        )
        reply = response.choices[0].message.content or ""
        return reply.replace("```json", "").replace("```", "")
    except Exception as e:
        print(f"⚠️ Cleanup LLM error: {e}")
        return '{"impacts": []}'

# =======================================
# 🔧 Helpers
# =======================================
def safe_json_load(s: str) -> Dict[str, Any]:
    try:
        return json.loads(s)
    except json.JSONDecodeError:
        return {"impacts": []}

def expand_impacts_to_rows(data: Dict[str, Any], source_file: str) -> List[Dict[str, Any]]:
    rows = []
    for impact in data.get("impacts", []):
        impact_text = (impact.get("text") or "").strip()
        industries = impact.get("Industry", []) or [{"value": "", "confidence": None}]
        business_groups = impact.get("Business Group", []) or [{"value": "", "confidence": None}]
        use_cases = impact.get("Use Case", []) or [{"value": "", "confidence": None}]

        for ind in industries:
            for bg in business_groups:
                for uc in use_cases:
                    rows.append({
                        "source_file": os.path.basename(source_file),
                        "impact": impact_text,
                        "industry": ind.get("value", ""),
                        "industry_conf": ind.get("confidence", None),
                        "business_group": bg.get("value", ""),
                        "business_group_conf": bg.get("confidence", None),
                        "use_case": uc.get("value", ""),
                        "use_case_conf": uc.get("confidence", None),
                        "source_type": SOURCE_TYPE,
                        "slide_no": None,  # basic extractor doesn't track slide_no
                    })
    return rows

def compute_fingerprint(row: Dict[str, Any]) -> str:
    sig = "|".join([
        (row.get("impact") or "").strip().lower(),
        (row.get("industry") or "").strip().lower(),
        (row.get("business_group") or "").strip().lower(),
        (row.get("use_case") or "").strip().lower(),
        (row.get("source_file") or "").strip().lower(),
        (row.get("source_type") or "").strip().lower(),
    ])
    return hashlib.sha256(sig.encode("utf-8")).hexdigest()

def pick_top_n_per_deck(df: pd.DataFrame, n: int) -> pd.DataFrame:
    if df.empty:
        return df
    df = df.copy()
    # average confidence across the three (NaN safe)
    df["avg_conf"] = pd.to_numeric(df["industry_conf"], errors="coerce").fillna(0)
    df["avg_conf"] += pd.to_numeric(df["business_group_conf"], errors="coerce").fillna(0)
    df["avg_conf"] += pd.to_numeric(df["use_case_conf"], errors="coerce").fillna(0)
    df["avg_conf"] = df["avg_conf"] / 3.0

    # enforce min confidence across all three dims
    df = df[
        (pd.to_numeric(df["industry_conf"], errors="coerce") >= MIN_CONF) &
        (pd.to_numeric(df["business_group_conf"], errors="coerce") >= MIN_CONF) &
        (pd.to_numeric(df["use_case_conf"], errors="coerce") >= MIN_CONF)
    ]

    if df.empty:
        return df

    # top-N per deck (source_file)
    df = df.sort_values(["source_file", "avg_conf"], ascending=[True, False])
    df = df.groupby("source_file", as_index=False).head(n)
    # compute fingerprint for dedupe
    df["fingerprint"] = df.apply(lambda r: compute_fingerprint(r.to_dict()), axis=1)
    df = df.drop_duplicates(subset=["fingerprint"])
    return df

def ensure_table(engine) -> None:
    with engine.begin() as conn:
        for stmt in DDL_CREATE.strip().split(";\n"):
            if stmt.strip():
                conn.execute(text(stmt))

def upsert_rows(engine, df: pd.DataFrame):
    """
    Upsert PPT impact data into PostgreSQL using SQLAlchemy + pgvector.
    Drops any extra columns not in the DB schema automatically.
    """
    from sqlalchemy import MetaData, Table
    from sqlalchemy.dialects.postgresql import insert as pg_insert

    TABLE_NAME = "impact_repository"
    SCHEMA_NAME = "scout"

    metadata = MetaData(schema=SCHEMA_NAME)
    table = Table(TABLE_NAME, metadata, autoload_with=engine)

    # Filter DataFrame columns to only those in table
    valid_cols = set(table.columns.keys())
    df_filtered = df[[c for c in df.columns if c in valid_cols]].copy()

    # Convert to records
    records = df_filtered.to_dict(orient="records")

    with engine.begin() as conn:
        stmt = pg_insert(table).values(records)
        stmt = stmt.on_conflict_do_nothing(index_elements=["fingerprint"])
        result = conn.execute(stmt)

    print(f"✅ {result.rowcount} records inserted into {SCHEMA_NAME}.{TABLE_NAME}")
    return result.rowcount



# =======================================
# 🚀 Main Async Flow
# =======================================
async def main(folder: Optional[str] = None) -> pd.DataFrame:
    folder = folder or FOLDER_PATH
    all_rows: List[Dict[str, Any]] = []

    ppt_files = [
        os.path.join(folder, f)
        for f in os.listdir(folder)
        if f.lower().endswith((".pptx", ".ppt"))
    ]

    if not ppt_files:
        print(f"ℹ️ No PPT/PPTX found under: {folder}")
        return pd.DataFrame()

    for ppt_path in tqdm(ppt_files, desc="Ingesting PPTs"):
        print(f"\n📊 Processing {os.path.basename(ppt_path):<80}")

        # Step 1: Extract text
        ppt_text = extract_text_from_ppt(ppt_path)
        if not ppt_text.strip():
            print(f"🟡 No text extracted from {ppt_path}. (Consider robust extractor w/ tables+notes)")
            continue

        # Step 2: Get LLM classification
        raw_json_str = await analyze_ppt_text(ppt_text)

        # Step 3: Clean irrelevant entries
        clean_json_str = await cleanup_irrelevant_impacts(raw_json_str)

        # Step 4: Parse JSON safely
        data = safe_json_load(clean_json_str)

        # Step 5: Expand to rows
        rows = expand_impacts_to_rows(data, ppt_path)
        if not rows:
            print(f"ℹ️ No extracted impacts for {ppt_path}")
        all_rows.extend(rows)

    # Step 6: Build DataFrame
    df_expanded = pd.DataFrame(all_rows,
        columns=[
            "source_file","impact","industry","industry_conf",
            "business_group","business_group_conf",
            "use_case","use_case_conf","source_type","slide_no"
        ]
    )

    if df_expanded.empty:
        print("\nℹ️ Nothing to load. Exiting.")
        return df_expanded

    # Step 7: Top-N per deck + min conf + fingerprint
    df_final = pick_top_n_per_deck(df_expanded, TOP_N_PER_DECK)

    if df_final.empty:
        print(f"\nℹ️ After filtering (min_conf={MIN_CONF}, top_n={TOP_N_PER_DECK}), nothing to load.")
        return df_final

    # Step 8: DB write (ensure + upsert)
    engine = get_engine()
    ensure_table(engine)
    inserted = upsert_rows(engine, df_final)

    # Report
    print(f"\n✅ PPT ingest complete. Inserted={inserted} rows "
          f"(post-filter unique={len(df_final)}; raw_total={len(df_expanded)})")

    return df_final

# =======================================
# 🏁 CLI Entrypoint
# =======================================
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Ingest PPT impacts into Azure Postgres (scout.impact_repository)")
    parser.add_argument("folder", nargs="?", default=FOLDER_PATH, help="Folder with PPT/PPTX files")
    parser.add_argument("--min-conf", type=float, default=MIN_CONF, help="Min confidence threshold (all three dims)")
    parser.add_argument("--top-n", type=int, default=TOP_N_PER_DECK, help="Top-N impacts per deck to keep")
    args = parser.parse_args()

    # Allow overriding via CLI
    MIN_CONF = args.min_conf
    TOP_N_PER_DECK = args.top_n
    FOLDER_PATH = args.folder

    # Run
    df_out = asyncio.run(main(FOLDER_PATH))
