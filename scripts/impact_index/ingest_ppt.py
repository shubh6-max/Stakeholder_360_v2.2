"""
Local PPTs -> pgvector (scout.cs_meta / scout.cs_chunks)
Content-first: derive metadata from slide content, not filenames.
"""

import os, hashlib
from time import perf_counter
from typing import List
from sqlalchemy.orm import Session
from pptx import Presentation

from s360_rag.config import assert_required, PPT_FOLDER, RAG_CHUNK_TOKENS, RAG_CHUNK_OVERLAP
from s360_rag.db import engine, SessionLocal, create_tables, CsMeta, CsChunks
from s360_rag.retriever import embed_text

def _read_ppt_text(path: str) -> List[str]:
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        joined = "\n".join([t.strip() for t in texts if t and t.strip()])
        slides.append((i, joined))
    return slides

def _split_tokens(text: str, max_tokens: int, overlap: int) -> List[str]:
    # simple char-based splitter as a placeholder (tokenizer-free)
    # you can replace with langchain_text_splitters RecursiveCharacterTextSplitter
    max_chars = max_tokens * 4  # rough approx
    ov_chars = overlap * 4
    chunks = []
    i = 0
    while i < len(text):
        chunk = text[i : i + max_chars]
        chunks.append(chunk)
        i += max_chars - ov_chars
    return chunks

def run():
    assert_required()
    create_tables()

    files = [os.path.join(PPT_FOLDER, f) for f in os.listdir(PPT_FOLDER) if f.lower().endswith(".pptx")]
    if not files:
        print(f"⚠️ No PPTX files found in {PPT_FOLDER}")
        return

    with SessionLocal() as session:
        for fp in files:
            print(f"🔹 Ingest: {os.path.basename(fp)}")
            slides = _read_ppt_text(fp)
            all_text = "\n\n".join([t for _, t in slides if t])
            if not all_text.strip():
                print("   (skip: empty content)")
                continue

            # Derive meta from content
            first_title = slides[0][1].splitlines()[0] if slides and slides[0][1] else "Untitled"
            case_id = hashlib.sha1((fp + first_title).encode("utf-8")).hexdigest()

            # Upsert cs_meta (simple delete-insert for idempotence)
            session.query(CsMeta).filter(CsMeta.case_id == case_id).delete()
            session.add(CsMeta(
                case_id=case_id,
                case_title=first_title[:200],
                client=None, industry=None, business_function=None,
                source_file=os.path.basename(fp),
                total_pages=len(slides)
            ))
            session.commit()

            # Chunk & embed per slide
            session.query(CsChunks).filter(CsChunks.case_id == case_id).delete()
            batch = []
            for page_no, text in slides:
                for chunk in _split_tokens(text, RAG_CHUNK_TOKENS, RAG_CHUNK_OVERLAP):
                    if not chunk.strip():
                        continue
                    emb = embed_text(chunk)
                    batch.append(CsChunks(case_id=case_id, page_no=page_no, chunk=chunk, embedding=emb))

            # bulk insert
            session.bulk_save_objects(batch)
            session.commit()
            print(f"   ✅ {len(batch)} chunks inserted")

if __name__ == "__main__":
    run()