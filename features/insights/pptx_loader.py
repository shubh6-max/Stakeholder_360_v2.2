# features/insights/case_parsers/pptx_loader.py
from __future__ import annotations

import os
import re
from typing import List, Tuple, Dict, Any

from langchain_core.documents import Document
from pptx import Presentation


def _clean_text(s: str) -> str:
    if not s:
        return ""
    s = s.replace("\x00", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\u00AD", "", s)  # soft hyphen
    s = re.sub(r"[ \t]*\n[ \t]*", "\n", s)
    return s.strip()


def _slide_text(slide) -> str:
    parts: List[str] = []
    # Shapes (title, body, tables, etc.)
    for shp in slide.shapes:
        # text frames
        if hasattr(shp, "text") and isinstance(shp.text, str):
            txt = _clean_text(shp.text)
            if txt:
                parts.append(txt)
        # tables
        if getattr(shp, "has_table", False):
            tbl = shp.table
            rows = []
            for r in tbl.rows:
                cells = [_clean_text(c.text) for c in r.cells]
                rows.append(" | ".join(c for c in cells if c))
            if rows:
                parts.append("\n".join(rows))
    return "\n".join([p for p in parts if p]).strip()


def load_pptx_as_documents(
    path: str,
    *,
    include_empty: bool = False,
) -> Tuple[List[Document], str, List[Dict[str, Any]]]:
    """
    Load a PPTX into LangChain Documents (one per slide).

    Returns:
        docs: List[Document] with page_content per slide
        full_text: concatenated text of all slides
        slide_page_text: [{"index": i, "text": "..."}]
    """
    if not os.path.exists(path):
        raise FileNotFoundError(path)

    prs = Presentation(path)
    docs: List[Document] = []
    slide_page_text: List[Dict[str, Any]] = []

    file_name = os.path.basename(path)
    source_path = os.path.abspath(path)

    for i, slide in enumerate(prs.slides):
        txt = _slide_text(slide)
        if include_empty or txt.strip():
            meta = {
                "file_name": file_name,
                "source_path": source_path,
                "source_type": "pptx",
                "slide_index": i,
                "page_index": i,  # symmetry with PDF
            }
            docs.append(Document(page_content=txt, metadata=meta))
            slide_page_text.append({"index": i, "text": txt})

    # ✅ FIX: build full_text correctly from slide_page_text
    full_text = _clean_text("\n\n".join(x["text"] for x in slide_page_text))

    return docs, full_text, slide_page_text